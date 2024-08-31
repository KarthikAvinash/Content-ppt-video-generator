import json
from PIL import ImageFont, ImageDraw, Image
import os
import google.generativeai as genai
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import requests
from pptx import Presentation
from pptx.util import Inches
import os
import requests
from pptx.dml.color import RGBColor
from pptx.util import Pt

if not load_dotenv('./.env'): raise Exception("env file not found")
GEMINI_API_KEY = os.getenv("GOOGLE_GEMINI_API_KEY")
genai.configure(api_key=GEMINI_API_KEY)
MODEL = "gemini-1.5-flash-latest"
# MODEL = "gemini-1.5-pro-latest"
model = genai.GenerativeModel(MODEL)

ppt_file = ""
fonts_dir = "fonts"

def get_text_dimensions(text, font):
    image = Image.new('RGB', (1, 1))
    draw = ImageDraw.Draw(image)
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    return text_width, text_height

def estimate_word_count(width_px, height_px, font_path, font_size):
    font = ImageFont.truetype(font_path, font_size)
    sample_text = "The quick brown fox jumps over the lazy dog"
    text_width, text_height = get_text_dimensions(sample_text, font)
    chars_per_line = width_px / text_width
    lines_per_box = height_px / text_height
    words_per_line = len(sample_text.split())
    total_words = words_per_line * lines_per_box
    return total_words

def extract_slide_info(slide):
    slide_info = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            width_px = width / 9525
            height_px = height / 9525
            if text in ['img', 'title', 'txt','stitle']:
                slide_info.append({
                    'placeholder': text,
                    'left': left,
                    'top': top,
                    'width': width_px,
                    'height': height_px
                })
    return slide_info

def create_prompt(slide_info):
    prompts = []
    for info in slide_info:
        placeholder = info['placeholder']
        left = info['left']
        top = info['top']
        width = info['width']
        height = info['height']
        prompt_entry = {
            'placeholder': placeholder,
            'location': {
                'left': left,
                'top': top
            },
            # 'dimensions': {
            #     'width': width,
            #     'height': height
            # }
        }
        if placeholder == 'img':
            prompt_entry["estimated words"] = f"10"
            prompt_entry['instructions'] = "[Write Image Here]"
        elif placeholder == 'title':
            font_path = r'fonts/Arial_Bold.ttf'
            font_size = 24
            estimated_words = estimate_word_count(width, height, font_path, font_size)
            prompt_entry["estimated words"] = f"{estimated_words:.0f}"
            prompt_entry['instructions'] = f"[Write Title Here]"
        elif placeholder == 'txt':
            font_path = r'fonts/Arial.ttf'
            font_size = 12
            estimated_words = estimate_word_count(width, height, font_path, font_size)
            prompt_entry["estimated words"] = f"{estimated_words:.0f}"
            prompt_entry['instructions'] = f"[Write Text Here]"
        elif placeholder=="stitle":
            font_path = r'fonts/Arial.ttf'
            font_size = 12
            estimated_words = estimate_word_count(width, height, font_path, font_size)
            prompt_entry["estimated words"] = f"{estimated_words:.0f}"
            prompt_entry['instructions'] = f"[Write Sub-Title Here]"
        prompts.append(prompt_entry)
    return prompts

def preprocess_response(text):
    text = text.replace('*', '')
    return text

def parse_response_file(file_path):
    instructions = {}
    
    with open(file_path, 'r') as file:
        lines = file.readlines()
        instructions = []
        for line in lines:
            if "\"instructions\":" in line:
                # Remove 'instructions: ' part
                # print(line)
                clean_instruction = line.split("\"instructions\": ")[1]

                print(clean_instruction[1:-2])
                instructions.append(clean_instruction[1:-2])
    print("instructions: ", instructions)
    return instructions


def update_presentation_with_instructions(pptx_path, instructions, output_path):
    presentation = Presentation(pptx_path)
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in instructions:
                    # Update shape text with instructions
                    shape.text_frame.clear()
                    p = shape.text_frame.add_paragraph()
                    p.text = instructions[text]
    
    presentation.save(output_path)
    print(f"Updated presentation saved as '{output_path}'.")

def remove_image_placeholder(instructions):
    updated_instructions = []
    for instruction in instructions:
        # Remove the "Image Placeholder" part from the string
        cleaned_instruction = instruction.replace("Image Placeholder: ", "")
        updated_instructions.append(cleaned_instruction)
    return updated_instructions

def emus_to_pixels(emus, dpi=96):
    # Convert EMUs to inches
    inches = emus / 914400
    # Convert inches to pixels
    pixels = inches * dpi
    return pixels

def fetch_and_save_image(query, width, height, save_dir, image_name, access_key):
    # Create the directory if it doesn't exist
    if not os.path.exists(save_dir):
        os.makedirs(save_dir)

    # Search for photos on Unsplash
    response = requests.get(
        "https://api.unsplash.com/search/photos",
        params={
            "query": query,
            "page": 1,
            "per_page": 1,  # Get only the first image
        },
        headers={"Authorization": f"Client-ID {access_key}"},
    )

    # Parse the response
    data = response.json()
    if data["results"]:
        # Get the raw URL of the first image
        raw_url = data["results"][0]["urls"]["raw"]
        
        # Append custom width and height to the URL
        custom_url = f"{raw_url}&w={width}&h={height}"
        
        # Fetch the image data
        image_response = requests.get(custom_url)

        if image_response.status_code == 200:
            # Define the path to save the image
            save_path = os.path.join(save_dir, f"{image_name}.jpg")
            
            # Save the image
            with open(save_path, 'wb') as file:
                file.write(image_response.content)
            
            return save_path  # Return the path to the saved image
        else:
            print("Failed to fetch the image.")
            return None
    else:
        print("No images found.")
        return None

def replace_placeholders_with_text(presentation_path, output_path, instructions, access_key):
    # Load the PowerPoint presentation
    presentation = Presentation(presentation_path)
    instr_idx = 0

    # Iterate over each slide in the presentation
    for slide in presentation.slides:
        # Iterate over each shape in the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                # print("width: ", width,"; height: ", height)

                for paragraph in text_frame.paragraphs:
                    curr = paragraph.text.strip()
                    if curr in ['img', 'title', 'txt', 'stitle']:
                        # Clear existing text
                        p = paragraph
                        p.clear()

                        if instr_idx < len(instructions):
                            new_text = instructions[instr_idx].strip('"')
                            instr_idx += 1

                            if curr == "img":
                                curr_width_pix = emus_to_pixels(width)
                                curr_height_pix = emus_to_pixels(height)
                                # Fetch and insert the image
                                image_path = fetch_and_save_image(
                                    query=new_text,
                                    width=curr_width_pix,  # You can adjust the width and height as needed
                                    height=curr_height_pix,
                                    save_dir="online_images",
                                    image_name=f"image_{instr_idx}",
                                    access_key=access_key
                                )
                                if image_path:
                                    # Replace shape with the image
                                    shape.element.getparent().remove(shape.element)
                                    slide.shapes.add_picture(image_path, left, top, width, height)  # Adjust size and position
                            else:
                                run = p.add_run()
                                run.text = new_text
                                run.font.color.rgb = RGBColor(255, 255, 255)

                                if curr == "title":
                                    run.font.bold = True
                                    run.font.size = Pt(40)
                                elif curr == "stitle":
                                    run.font.size = Pt(25)
                                else:
                                    run.font.size = Pt(16)

    # Save the updated presentation
    presentation.save(output_path)
    print(f"Updated PowerPoint file saved as '{output_path}'.")