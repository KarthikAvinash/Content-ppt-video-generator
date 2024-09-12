import streamlit as st
import os
from pptx import Presentation
import json
from utils.ppt_util_functions import *
from utils.video_util_functions import *
import io
from PIL import Image
import tempfile
from pptxtopdf import convert
# New imports
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import re
import json
import os

image_dir = "images"

# Function to save uploaded file
def save_uploaded_file(uploaded_file, save_path):
    with open(save_path, "wb") as f:
        f.write(uploaded_file.getbuffer())

# New function to convert PowerPoint to images
def pptx_to_images(pptx_path):
    prs = Presentation(pptx_path)
    images = []
    for slide in prs.slides:
        img = slide_to_image(slide)
        images.append(img)
    return images

def slide_to_image(slide):
    image = Image.new("RGB", (1920, 1080))  # Adjust size as needed
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if shape.placeholder_format.type == 1:  # Title
                title = shape.text_frame.text
                # Draw title on image (you may need to adjust font, size, position)
        # Add more conditions to handle other shape types
    return image

# Function to get sorted image files from a directory
def get_sorted_images(directory):
    image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
    image_files = [f for f in os.listdir(directory) if os.path.splitext(f.lower())[1] in image_extensions]
    return sorted(image_files)

def save_pptx_copy(original_path, new_path):
    """
    Reads a PowerPoint (.pptx) file and saves a copy with a different name.
    
    Parameters:
        original_path (str): The path to the original PowerPoint file.
        new_path (str): The path where the new PowerPoint file will be saved.
    """
    try:
        # Load the original presentation
        presentation = Presentation(original_path)
        
        # Save a copy of the presentation
        presentation.save(new_path)
        
        print(f"Presentation saved successfully to {new_path}")
    except Exception as e:
        print(f"Failed to save presentation: {e}")

def parse_response_file_to_get_num_instructions(filepath):
    instructions_count = {}
    slide_num_pattern = re.compile(r'"Slide (\d+)": \[')  # Regex to capture Slide <number> pattern
    instructions_pattern = re.compile(r'"instructions"')  # Pattern to identify 'instructions'

    with open(filepath, 'r', encoding='ISO-8859-1') as file:
        content = file.readlines()

    current_slide = None
    instruction_lines = 0

    for line in content:
        # Check for the start of a new slide
        slide_match = slide_num_pattern.search(line)
        if slide_match:
            # Save the count for the previous slide if applicable
            if current_slide is not None:
                instructions_count[current_slide] = instruction_lines

            # Reset for the new slide
            current_slide = slide_match.group(1)
            instruction_lines = 0  # Reset instruction count for the new slide

        # Check if the line contains the word 'instructions'
        if instructions_pattern.search(line):
            instruction_lines += 1

    # For the last slide after the loop ends
    if current_slide is not None:
        instructions_count[current_slide] = instruction_lines

    return instructions_count


INSTRUCTIONS_FILE = "instructions.txt"
RESPONSE_FILE = 'response.txt'
def save_instructions(instructions):
    """Save instructions to the text file."""
    with open(INSTRUCTIONS_FILE, "w") as f:
        f.write("\n".join(instructions))


JSON_FILE = "instruction_count.json"

def save_json(data, filename=JSON_FILE):
    """Saves a dictionary to a JSON file."""
    with open(filename, 'w') as json_file:
        json.dump(data, json_file, indent=4)
    print(f"Data successfully saved to {filename}")

def load_json(filename=JSON_FILE):
    """Loads a dictionary from a JSON file."""
    if os.path.exists(filename):
        with open(filename, 'r') as json_file:
            data = json.load(json_file)
        print(f"Data successfully loaded from {filename}")
        return data
    else:
        print(f"No file found at {filename}, returning empty dictionary.")
        return {}


import requests
from bs4 import BeautifulSoup as bs

def fetch_image_url(query):
    """Fetch the URL of an image based on the search query."""
    params = {
        "q": query,  # Search query
        "tbm": "isch",  # Image results
    }
    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
    }
    try:
        html = requests.get("https://www.google.com/search", params=params, headers=headers, timeout=30)
        soup = bs(html.content, "html.parser")
        images = soup.select('div img')
        if images and 'src' in images[1].attrs:
            return images[1]['src']
    except Exception as e:
        print(f"Error fetching image URL: {e}")
    return ""

#_____________________________

ppt_dir = "ppts"
ppt_path = os.path.join(ppt_dir, "template.pptx")

def get_ppt():
    st.title("PowerPoint Template Updater")

    # File upload section
    col1, col2 = st.columns(2)
    with col1:
        template_file = st.file_uploader("Upload PowerPoint Template", type=["pptx"])
    with col2:
        content_file = st.file_uploader("Upload Content File", type=["txt"])

    if template_file and content_file:
        # Save uploaded files
        
        os.makedirs(ppt_dir, exist_ok=True)
        
        content_path = "content.txt"

        save_uploaded_file(template_file, ppt_path)
        save_uploaded_file(content_file, content_path)

        # Textbox for user to specify the topic
        user_topic = st.text_input("Enter the topic you want to create the presentation on: (Type 'NA' if you do not have any specific topic and wish to use the whole content file for ppt generation)")
        add_to_prompt = ""
        if user_topic:
            st.write(f"Topic entered: {user_topic}")
            print(f"Topic entered by user: {user_topic}")
            add_to_prompt = " and this is request on which the ppt is to be made - "+user_topic + "."

        # Process files
        if st.button("Process Files"):
            progress_bar = st.progress(0)
            status_text = st.empty()

            # Step 1: Extract slide info and create prompts
            status_text.text("Extracting slide info and creating prompts...")
            presentation = Presentation(ppt_path)
            all_prompts = {}
            for slide_num, slide in enumerate(presentation.slides, start=1):
                slide_info = extract_slide_info(slide)
                slide_prompts = create_prompt(slide_info)
                all_prompts[f"Slide {slide_num}"] = slide_prompts
            progress_bar.progress(25)

            # Step 2: Save prompts to JSON
            status_text.text("Saving prompts to JSON...")
            prompts_dir = "prompts"
            os.makedirs(prompts_dir, exist_ok=True)
            prompt_file = os.path.join(prompts_dir, "prompt_slides_info.json")
            with open(prompt_file, 'w') as f:
                json.dump(all_prompts, f, indent=4)
            progress_bar.progress(50)

            # Step 3: Generate content
            status_text.text("Generating content...")
            with open(content_path, 'r', encoding="utf-8") as content_file:
                content = content_file.read()
            with open(prompt_file) as query_file:
                query = query_file.read()
            

            fin_query = f"I need to create a presentation {add_to_prompt}. I need assistance in filling the placeholders in the given template which is in the form of a json. You need to only fill the 'instructions' section in each one in each slide. And return that json content only, so that I can put it directly in a json file. The presentation is based on the following content: {content}\n\nand the template you need to fill using the given above content is as follows, make sure that the word limit is taken care and keep it as short as possible, and also the content in each slide must be linked, and all slides must be linked, and all content must be related to the given content. Write the points in short and precise. Also forget the image placeholders in this template: {query}"
            
            response = model.generate_content([fin_query], stream=True)
            response.resolve()
            output = preprocess_response(response.text)
            
            with open("response.txt", "w") as file:
                file.write(output)
            progress_bar.progress(75)

            # Step 4: Update presentation
            status_text.text("Updating presentation...")
            instructions = parse_response_file(RESPONSE_FILE)
            instruction_count = parse_response_file_to_get_num_instructions(RESPONSE_FILE)  # E.g., {'1': 3, '2': 11, '3': 9}
            save_instructions(instructions)
            save_json(instruction_count)

# copied to edit_instructions.py as well.

            update_presentation_with_instructions(ppt_path, instructions, 'updated_presentation.pptx')
            instructions = remove_image_placeholder(instructions)

            output_path = 'updated_presentation.pptx'
            access_key = "IOuG8tW-MhjIxRKKXrhd079mIqGArt9iqF4NLn1AFsE"
            replace_placeholders_with_text(ppt_path, output_path, instructions, access_key)
            progress_bar.progress(100)

            status_text.text("Processing complete!")
            
            # Provide download option for updated presentation
            with open(output_path, "rb") as file:
                st.download_button(
                    label="Download Updated Presentation",
                    data=file,
                    file_name="updated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
